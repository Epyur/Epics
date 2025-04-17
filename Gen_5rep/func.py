import ast
import itertools
import logging
import os
from datetime import datetime
from pathlib import Path
import openpyxl
import pandas as pd
import numpy as np
from matplotlib import pyplot as plt
from mpl_toolkits.mplot3d import Axes3D
from openpyxl import *
from openpyxl.chart import ScatterChart, Reference
from openpyxl.chart.series import Series
from openpyxl.drawing.image import Image
from openpyxl.reader.excel import load_workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

logging.basicConfig(level=logging.ERROR)


def find_matching_excel_files(root_dir, exclude_dirs=None):
    """
    Находит все xlsx файлы, имена которых совпадают с названием содержащей их папки.

    Параметры:
    - root_dir: корневая директория для поиска
    - exclude_dirs: список директорий, которые нужно исключить из поиска

    Возвращает:
    Список путей к найденным файлам
    """

    if exclude_dirs is None:
        exclude_dirs = ['Архив']

    matching_files = []

    for dirpath, dirnames, filenames in os.walk(root_dir):
        # Исключаем заданные директории из обхода
        dirnames[:] = [d for d in dirnames if d not in exclude_dirs]

        for dirname in dirnames:
            # Формируем полное имя файла, которое должно совпадать с названием папки
            expected_file = os.path.join(dirpath, dirname, f"{dirname}.xlsx")

            # Проверяем, существует ли такой файл
            if os.path.exists(expected_file):
                matching_files.append(expected_file)

    return matching_files

def DfList(initial_rout):
    files = find_matching_excel_files(initial_rout)
    dataframes = []
    for file in files:
        df = pd.read_excel(file, index_col=0)
        df_transposed = df
        dataframes.append(df_transposed)
    return dataframes

def MergDF(df_list):
    try:
        # Проверка входных данных
        if not isinstance(df_list, list):
            raise TypeError("df_list должен быть списком датафреймов")

        # Объединение датафреймов
        merged_df = pd.concat(df_list, ignore_index=True)
        # Установка и сортировка индекса
        s_df = merged_df.reset_index(drop=True, level=0)

    except FileNotFoundError as e:
        print(f"Ошибка: файл не найден - {e}")
    except ValueError as e:
        print(f"Ошибка: {e}")
    except TypeError as e:
        print(f"Ошибка: неверный тип данных - {e}")
    except Exception as e:
        print(f"Произошла ошибка: {e}")
    return s_df


def ColorMark(file_nam, red_key, green_key):
    try:
        # Проверка существования файла
        if not Path(file_nam).exists():
            raise FileNotFoundError(f"Файл {file_nam} не найден")

        # Читаем существующий файл
        existing_df = pd.read_excel(file_nam, sheet_name=None)

        # Создаем writer для перезаписи файла
        writer = pd.ExcelWriter(file_nam, engine='xlsxwriter')
        workbook = writer.book

        # Форматирование для красного цвета
        format_red = workbook.add_format({
            'bg_color': '#FFC7CE',
            'font_color': '#9C0006'
        })

        # Форматирование для зеленого цвета
        format_green = workbook.add_format({
            'bg_color': '#00FF00',
            'font_color': '#000000'
        })

        # Формат для заголовков
        header_format = workbook.add_format({
            'bg_color': '#F2F2F2',
            'text_wrap': True,
            'align': 'center',
            'valign': 'vcenter',
            'border': 1
        })

        # Общий формат ячеек
        cell_format = workbook.add_format({
            'border': 1,
            'text_wrap': True,
            'align': 'center',
            'valign': 'vcenter'
        })

        # Проходим по всем листам
        for sheet_name, df in existing_df.items():
            try:
                # Добавляем датафрейм на лист
                df.to_excel(writer, sheet_name=sheet_name, index=False, header=True)

                # Получаем лист
                worksheet = writer.sheets[sheet_name]

                # Установка области форматирования
                last_row = len(df) + 1
                format_range = f'A1:CT{last_row}'

                # Условие для красного цвета
                worksheet.conditional_format(
                    format_range,
                    {
                        'type': 'text',
                        'criteria': 'containing',
                        'value': red_key,
                        'format': format_red
                    }
                )

                # Условие для зеленого цвета
                worksheet.conditional_format(
                    format_range,
                    {
                        'type': 'text',
                        'criteria': 'containing',
                        'value': green_key,
                        'format': format_green
                    }
                )

                # Настройка заголовков
                worksheet.set_row(0, 40)
                worksheet.write_row(0, 0, df.columns, header_format)


                # Автоматический подбор ширины колонок
                for i, col in enumerate(df.columns):
                    column_len = max(df[col].astype(str).str.len().max(), len(col))
                    worksheet.set_column(i, i, column_len * 1.1)

                # Применить формат к данным
                worksheet.conditional_format(
                    'A2:' + format_range.split(':')[1],
                    {
                        'type': 'no_errors',
                        'format': cell_format
                    }
                )

                print(f"Форматирование успешно применено к листу {sheet_name}")
            except Exception as e:
                print(f"Ошибка при обработке листа {sheet_name}: {e}")

        # Сохраняем изменения
        writer.close()
        print(f"Форматирование успешно применено ко всем листам файла {file_nam}")

    except FileNotFoundError as e:
        print(f"Ошибка: {e}")
    except Exception as e:
        print(f"Произошла ошибка: {e}")


def SelectAndFilter(d_f,
                    select_column,
                    mean_column = None,
                    key_column = None,
                    key_word = None,
                    column_name=None,
                    filters=None,
                    date_column=None,
                    file_nam=None,
                    sheet_nam=None,
                    dict_to_rename=None,
                    start_date=None,
                    end_date=None,
                    quarter=None,
                    half_year=None,
                    drop_date=None):
    try:
        # Проверка на None перед работой с датафреймом
        if d_f is None:
            raise ValueError("Входной датафрейм не должен быть None")

        filtered_data = d_f.reset_index()
        filtered_data = filtered_data[select_column]
        if mean_column:
            filtered_data = filtered_data[filtered_data[mean_column].notna()]
        if key_word:
            filtered_data = filtered_data[filtered_data[key_column].str.contains(key_word, case=False, na=False)]
        else:
            pass

        if column_name:
            filtered_data = filtered_data[filtered_data[column_name].isin(filters)]

        # Конвертация столбца в формат даты
        if date_column:
            filtered_data[date_column] = pd.to_datetime(filtered_data[date_column], format='%d.%m.%Y', errors='coerce')

        # Фильтрация по дате
        if start_date:
            start_date = datetime.strptime(start_date, '%d.%m.%Y')
            filtered_data = filtered_data[filtered_data[date_column] >= start_date]

        if end_date:
            end_date = datetime.strptime(end_date, '%d.%m.%Y')
            filtered_data = filtered_data[filtered_data[date_column] <= end_date]

        # Фильтрация по кварталам
        if quarter:
            filtered_data = filtered_data[filtered_data[date_column].dt.quarter == quarter]

        # Фильтрация по полугодиям
        if half_year:
            if half_year == 1:
                filtered_data = filtered_data[filtered_data[date_column].dt.month <= 6]
            elif half_year == 2:
                filtered_data = filtered_data[filtered_data[date_column].dt.month > 6]
        if drop_date:
            filtered_data = filtered_data.drop(drop_date, axis=1)

        new_dataframe = filtered_data.copy()

        # Переименование столбцов
        if dict_to_rename is not None:
            new_dataframe = new_dataframe.rename(columns=dict_to_rename)

        # Сохранение в Excel, если требуется
        if file_nam:
            with pd.ExcelWriter(file_nam, mode='a', engine='openpyxl', if_sheet_exists='replace') as writer:
                new_dataframe.to_excel(writer, sheet_name=sheet_nam, index=False)

        return new_dataframe


    except Exception as e:
        print(f"Произошла ошибка: {e}")
        logging.exception("Произошла ошибка")


def plot_pie_chart(series, title='Круговая диаграмма', autopct='%.1f%%', figsize=(6, 6), colors=None, new_file_name=None, excel_file=None,
                   sheet_name='Sheet1', position=(1, 1), explode=None, perspective=0.5):
    """
    Построение объемной разрезной круговой диаграммы

    Параметры:
    - series: pandas Series с категориальными значениями
    - title: заголовок диаграммы
    - autopct: формат отображения процентов
    - figsize: размер фигуры
    - colors: список цветов для секторов
    - excel_file: путь к файлу Excel
    - sheet_name: имя листа в Excel
    - position: кортеж (row, col) для позиции вставки
    - explode: кортеж для разделения секторов (по умолчанию все 0)
    - perspective: угол перспективы (0-1)

    Возвращает:
    matplotlib axes object
    """
    # Подсчет значений
    counts = series.value_counts()

    # Если explode не указан, создаем кортеж с нулями
    if explode is None:
        explode = (0,) * len(counts)

    # Создание объемной диаграммы
    plt.figure(figsize=figsize)
    fig = plt.gcf()  # Получаем текущую фигуру
    ax = Axes3D(fig, azim=-60, elev=30, )  # Создаем 3D оси


    # Построение диаграммы
    wedgeprops = dict(width=perspective)
    patches, texts, autotexts = ax.pie(
        counts,
        explode=explode,
        autopct=autopct,
        colors=colors,
        startangle=90,
        labels=None,
        wedgeprops=wedgeprops,
        textprops={'fontsize': 12}
    )

    # Настройка внешнего вида
    ax.set_title(title, fontsize=16)

    # Создание списка меток для легенды
    labels = [f"{cat} ({pct:.1f}%)" for cat, pct in zip(counts.index, counts / counts.sum() * 100)]

    # Размещение легенды внизу по центру
    legend = plt.legend(
        patches,  # Используем патчи для легенды
        labels,
        title=series.name,
        loc='center',
        bbox_to_anchor=(0.5, -0.15),
        fontsize=12,
        ncol=len(counts)
    )
    # Сохранение графика в файл
    plt.savefig(f'{new_file_name}.jpg', bbox_inches='tight', dpi=300)

    return ax


def PicatorVtoroy(df, column_names, grafic_name, file_name, cat1, cat2, cat3=None,
                 subplot_order=None, category_order=None, width=12, hight=6):
    try:
        # Проверка входных данных
        if df.empty:
            raise ValueError("DataFrame is empty")
        if not isinstance(column_names, list):
            raise TypeError("column_names must be a list")

        # Устанавливаем порядок подграфиков
        if subplot_order is None:
            subplot_order = column_names
        else:
            column_names = subplot_order

        # Устанавливаем порядок категорий
        if cat3 is None:
            category_order = [cat1, cat2]
            # print(df.values.ravel())
        else:
            category_order = [cat1, cat2, cat3]

        # Получаем уникальные значения категорий
        unique_values = list(set(df.values.ravel()))

        # Создаем новый DataFrame с использованием декартова произведения
        result = pd.DataFrame({
            'column_names': [item for item in column_names for _ in range(len(unique_values))],
            'category': unique_values * len(column_names)
        })

        # Раскладываем множество значений по строкам
        result = result.explode('category')

        # Добавляем процентное содержание
        percentage_dict = {}
        for col in column_names:
            col_counts = df[col].value_counts(normalize=True).mul(100).round(2)
            percentage_dict[col] = col_counts

        result['percentage'] = result.apply(
            lambda row: percentage_dict[row['column_names']].get(row['category'], 0),
            axis=1
        )

        # Сортируем данные по заданному порядку
        result['category'] = pd.Categorical(result['category'],
                                          categories=category_order,
                                          ordered=True)
        result = result.sort_values(['column_names', 'category'])

        # Определяем цвета для каждой категории
        if cat3 is None:
            colors = {cat1: 'green', cat2: 'red'}
        else:
            colors = {cat1: 'green', cat2: 'red', cat3: 'gray'}
        # Создаем подграфики
        fig, axes = plt.subplots(nrows=1,
                                ncols=len(subplot_order),
                                figsize=(width, hight),
                                sharey=True)

        # Группируем данные
        grouped = result.groupby('column_names')

        # Создаем график для каждой группы
        for i, (name) in enumerate(subplot_order):
            group = grouped.get_group(name)
            # Рисуем бары с заданными цветами
            ax = group.plot.bar(x='category',
                               y='percentage',
                               ax=axes[i],
                               rot=90,
                               color=group['category'].map(colors))

            ax.yaxis.set_ticks(np.arange(0, 101, 5))

            # Добавляем название подграфика вертикально слева
            ax.set_title(name, rotation=90, loc='right', y=0, x=1, fontsize=9, fontweight='bold')

            # Добавляем подписи осей
            ax.set_xlabel('')
            ax.set_ylabel('Доля, %')
            ax.set_xticklabels('', rotation=0)


            # Добавляем сетку для лучшей читаемости
            ax.grid(axis='y', linestyle='--', alpha=0.7)

            # Удаляем индивидуальные легенды
            ax.legend_.remove()

        # Создаем общую легенду под графиками
        handles = [plt.Rectangle((0, 0), 1, 1, color=c) for c in colors.values()]
        legend = plt.legend(handles, colors.keys(), loc='lower right', bbox_to_anchor=(0.5, -0.1),
                            fancybox=True, shadow=False, ncol=len(colors), fontsize=10)

        # Общее название для всей фигуры
        fig.suptitle(grafic_name, fontsize=12, fontweight='bold')

        # # Убираем лишние пробелы между подграфиками
        # plt.tight_layout(pad=3.0)

        # Сохраняем график по предустановленному пути
        save_path = f'{file_name}.jpg'
        fig.savefig(save_path, bbox_inches='tight', dpi=300)
        print(f"График сохранен по пути: {save_path}")

        # plt.show()


    except Exception as e:
        print(f"Произошла ошибка: {e}")


def DateFormat(df, column_list):
    # Преобразуем столбец в формат datetime
    for i in column_list:
        df[i] = pd.to_datetime(df[i], format='%y-%m-%d')

        # Меняем формат на dd.mm.yyyy
        df[i] = df[i].dt.strftime('%d.%m.%Y')
    return df

def CheckPlaceholders(slide_name):
    for placeholder in slide_name.placeholders:
        print(placeholder.placeholder_format.idx, placeholder.placeholder_format.type)

def FormatText(place_holder, size=None, color=None, bold=None, italic=None, align=None):
    
    # Форматирование текста
    text_frame = place_holder.text_frame
    paragraph = text_frame.paragraphs[0]
    run = paragraph.runs[0]

    if size is not None:
        # Размер шрифта
        run.font.size = Pt(size)

    if color is not None:
        # Цвет шрифта
        run.font.color.rgb = RGBColor(0, 0, 255)  # синий цвет

    # Начертание
    if bold is not None:
        run.font.bold = True
    if italic is not None:    
        run.font.italic = True

    if align == 'L':
        # Выравнивание
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    if align == 'R':
        # Выравнивание
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
    if align == 'C':
        # Выравнивание
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    if align == 'J':
        # Выравнивание
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.JUSTIFY

# def TransparencyImg(shape, alpha):
    
#     ts = shape.fill._xPr.solidFill
#     sF = ts.get
#     sE = SubElement(sF, 'a:alpha', val=str(alpha))

def FirstTitleSlide(prs,
                    slide_name,
                    sub_name,
                    rout_to_pic, 
                    layout=2,                     
                    font_size=28, 
                    subtitleplace=13,
                    subsubtitleplace=12,
                    year='25',
                    img_place=11,
                    check_holders=False):
    slide_layout = prs.slide_layouts[layout]  # Заголовок и контент
    slide_title = prs.slides.add_slide(slide_layout)

    if check_holders is True:
        CheckPlaceholders(slide_title)
    else:
        pass

    title_first = slide_title.shapes.title
    title_first.text = slide_name
    title_first.text_frame.paragraphs[0].runs[0].font.size = Pt(font_size)
    subtitle_1 = slide_title.placeholders[subtitleplace]
    subtitle_1.text = sub_name
    subtitle_0 = slide_title.placeholders[subsubtitleplace]
    subtitle_0.text = year
    image_placeholder = slide_title.placeholders[img_place]
    image_placeholder.insert_picture(rout_to_pic)

def PartTitleSlide(prs,
                    slide_name,
                    sub_name,
                    rout_to_pic, 
                    layout=4,                    
                    titleplace=10,
                    subtitleplace=11,                    
                    img_place=12,
                    check_holders=False):
    slide_layout_part = prs.slide_layouts[layout]
    slide_part1_title = prs.slides.add_slide(slide_layout_part)
    
    if check_holders is True:
        CheckPlaceholders(slide_part1_title)

    title_first_part = slide_part1_title.placeholders[titleplace]
    title_first_part.text = slide_name
    subtitle_2 = slide_part1_title.placeholders[subtitleplace]
    subtitle_2.text = sub_name
    FormatText(title_first_part, align='C')
    image_placeholder_2 = slide_part1_title.placeholders[img_place]
    image_placeholder_2.insert_picture(rout_to_pic)

def JustTextSlide(prs,
                  slide_name,
                  text_list,
                  font_size=28,                  
                  layout=11,                    
                  textplace=13,                    
                  text_size=16,
                  R=0,
                  G=0,
                  B=0,
                  check_holders=False):
    # text_1 = ['1. Собрание СБЕ']
    slide_1_1_layouts = prs.slide_layouts[layout]
    slide_1_1 = prs.slides.add_slide(slide_1_1_layouts)
    if check_holders is True:
        CheckPlaceholders(slide_1_1)

    title_1_1 = slide_1_1.shapes.title
    title_1_1.text = slide_name
    title_1_1.text_frame.paragraphs[0].runs[0].font.size = Pt(font_size)
    body_1_1 = slide_1_1.placeholders[textplace]
    body_1_1_tf = body_1_1.text_frame
    body_1_1_tf.clear()
    count = -1
    count_list = []
    str_1 = ''
    for i in text_list:
        count += 1
        count_list.append(count)
        # body_1_1_tf.add_paragraph()
        str_1 = str_1 + f'{i} \n' 
        
    body_1_1_tf.text = str_1
    for x in range(0, len(count_list)):
        body_1_1_tf.paragraphs[x].runs[0].font.size = Pt(text_size)
        body_1_1_tf.paragraphs[x].runs[0].font.color.rgb = RGBColor(R, G, B)
    # body_1_1_pg = body_1_1_tf.add_paragraph()
    # body_1_1_pg.add_run().text = text_list

def FourSquareSlide(prs,
                  slide_name,
                  ruts,
                  sub_title_list=None,
                  pigs_list=None,
                  texts_list=None,
                  title_font_size=28,
                  subtitle_font_size=18,                  
                  layout=30,                                  
                  text_size=12,
                  R=0,
                  G=0,
                  B=0,
                  check_holders=False):
    # placeholders map
    if sub_title_list is None:
        sub_title_list = [13, 14, 16, 15]
    if pigs_list is None:
        pigs_list = [28, 22, 29, 30]
    if texts_list is None:
        texts_list = [26, 25, 27, 21]

    routs_names = []
    text_lists = []
    progress_list = []
    for i in ruts:
        routs_names.append(i[0])
        # print(i[0])
        text_lists.append(i[1])
        progress_list.append(i[2])
    subtitle_dict = dict(zip(sub_title_list, routs_names))
    # print(subtitle_dict)
    text_dict = dict(zip(texts_list, text_lists))
    progress_dict = dict(zip(pigs_list, progress_list))
    
      
    slide_layout_2_1 = prs.slide_layouts[layout]
    slide_2_1 = prs.slides.add_slide(slide_layout_2_1)

    if check_holders is True:
        CheckPlaceholders(slide_2_1)

    title_2_1 = slide_2_1.shapes.title
    title_2_1.text = slide_name
    title_2_1.text_frame.paragraphs[0].runs[0].font.size = Pt(title_font_size)

    for i in range(0, len(sub_title_list)):
        k = sub_title_list[i]
        v = routs_names[i]
        j = pigs_list[i]
        j_r = progress_list[i]
        t_i = texts_list[i]
        t = text_lists[i]
        item = slide_2_1.placeholders[k]
        item.text = v
        item.text_frame.paragraphs[0].runs[0].font.size = Pt(subtitle_font_size)
        image_placeholder_2_1_1 = slide_2_1.placeholders[j]
        if j_r==10:
            progress = r'pics\10.png'
        if j_r==25:
            progress = r'pics\25.png'
        if j_r==50:
            progress = r'pics\50.png'
        if j_r==75:
            progress = r'pics\75.png'
        if j_r==100:
            progress = r'pics\100.png'
        else:
            pass
        image_placeholder_2_1_1.insert_picture(progress)
        item_2_1_1_aim = slide_2_1.placeholders[t_i]
        item_2_1_1_aim = item_2_1_1_aim.text_frame
        str_1 = ''
        for a in t:
            str_1 = str_1 + f'{a} \n'
        item_2_1_1_aim.text = str_1
        for b in range(0, len(t)):
            item_2_1_1_aim.paragraphs[b].runs[0].font.size = Pt(text_size)
            item_2_1_1_aim.paragraphs[b].runs[0].font.color.rgb = RGBColor(R, G, B)

def ProjectSlide(prs,
                  slide_name,
                  text_list1,
                  text_list2,                  
                  font_size=28,                  
                  layout=32,                    
                  textplace1=14,                    
                  text_size=10,
                  textplace2=15,                    
                  text_size2=10,
                  img_place=13,
                  img_link=None,
                  R=0,
                  G=0,
                  B=0,
                  check_holders=False):
    slide_layout_2_2 = prs.slide_layouts[layout]
    slide_2_2 = prs.slides.add_slide(slide_layout_2_2)
    if check_holders is True:
        CheckPlaceholders(slide_2_2)

    image_placeholder_2_2_1 = slide_2_2.placeholders[img_place]
    image_placeholder_2_2_1.insert_picture(img_link)
    title_2_2 = slide_2_2.shapes.title
    title_2_2.text = slide_name
    title_2_2.text_frame.paragraphs[0].runs[0].font.size = Pt(font_size)
    item_2_2_aim = slide_2_2.placeholders[textplace1].text_frame
    str1 = 'Действия и намерения: \n'
    for txt in text_list1:
        str1 = str1 + f'{txt} \n'

    item_2_2_aim.text = str1
    item_2_2_aim.paragraphs[0].runs[0].font.size = Pt(18)
    item_2_2_aim.paragraphs[0].runs[0].font.color.rgb = RGBColor(22, 0, 102)
    for c_l in range(1, len(text_list1)+1):
        item_2_2_aim.paragraphs[c_l].runs[0].font.size = Pt(text_size)
        item_2_2_aim.paragraphs[c_l].runs[0].font.color.rgb = RGBColor(R, G, B)

    item_2_2_lim = slide_2_2.placeholders[textplace2]
    
    str2 = 'Ограничения: \n'
    for txt in text_list2:
        str2 = str2 + f'{txt} \n'
    item_2_2_lim.text = str2
    item_2_2_lim.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    item_2_2_lim.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(22, 0, 102)
    for c_e in range(1, len(text_list2)+1):
        item_2_2_lim.text_frame.paragraphs[c_e].runs[0].font.size = Pt(text_size2)
        item_2_2_lim.text_frame.paragraphs[c_e].runs[0].font.color.rgb = RGBColor(R, G, B)  