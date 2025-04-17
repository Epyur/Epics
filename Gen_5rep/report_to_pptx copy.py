import os
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from func import *

start_date = '01.01.2025'
end_date = '31.03.2025'
period_name = 'I квартал'

front_page_pic_rout = r'pics\title.png'
part1_title_pic_rout = r'pics\part1.png'
part2_title_pic_rout = r'pics\part2.png'
part3_title_pic_rout =r'pics\part3.jpeg'

report_file = r'rep\rep.pptx'

template = os.path.abspath(r'C:\Users\epyur\OneDrive\Документы\temple\temple.pptx')
prs = Presentation(template)
# prs = prs.slide_layouts[19]

'''Формируем титульный слайд'''

FirstTitleSlide(prs, 'Направление: "Пожарная безопасность"', period_name, front_page_pic_rout)

'''Закончили формировать титульный слайд'''

'''Начинаем 1 раздел: мероприятия'''

PartTitleSlide(prs, 'Техническая поддержка и мероприятия', period_name, part1_title_pic_rout)
text_list1 = ['18-20 февраля: Участие в собрании СБЕ;', 
    'январь - февраль: Участие в совещаниях по Саларьево (МИП);',
    'январь - февраль: Участие в проведении испытаний мембран в ЦЭИИС, в рамках работы по объекту "Саларьево";',
    'Общая техническая поддержка (подготовка писем, проведение консультаций);',
    'Получен сертификат на систему ТН-Кровля Гарант в системе НСОПБ с размещением информации о сертификате в открытых реестрах системы сертификации;',
    'Получен протокол Г1 на PIRCRYO, по результатам испытаний во ВНИИПО;',
    '30-31 января: Участие в конференции "Безопасная Арктика" в г. Оренбург;',
    '20 марта: Участие в конференции "Ройтмановские чтения" в АГПС (г. Москва) совместно с Жамойдиком С;',
    'Пожар в клину: выезд на объект (снятие сомнений по вопросу применения системы ТН-Кровля Гарант)']


JustTextSlide(prs, 'Текущая работа в I квартале 2025 года', text_list1)

'''Завершаем 1 раздел'''

'''Начинаем 2 раздел: НИОКР и проектная деятельность'''

PartTitleSlide(prs, 'НИОКРы и проектная деятельность', period_name, part2_title_pic_rout)

slide2_1_name = 'Основные направления движения'
rout_1 = 'Гарант в К0 (программа максимум)'
rout_1_v = ['Cнять внешние риски и ограничения для продвижения системы ТН-Кровля Гарант']
rout_1_p = 10
rout_2 = 'Гарант в К0 (программа минимум)'
rout_2_v = ['Снять внутрикорпоративные ограничения на продвижение системы ТН-Кровля Гарант']
rout_2_p = 25
rout_3 = 'Развитие пожарной комплектации'
rout_3_v = ['Развитие комплектности предложений СБЕ;', 'Формирование конкурентных преимуществ для системных решений СБЕ.']
rout_3_p = 75
rout_4 = 'Развитие ЛПИ и исследований материалов СБЕ'
rout_4_v = ['Провести модернизацию исследовательской базы лаборатории;',
            'Выстроить процесс контроля работы лаборатории, включая формирование базы результатов проводимых исследований;', 
            'Организовать проведение исследований показателей воспламеняемости и распространения пламени по поверхности ПВХ мембран.']
rout_4_p = 50

routs_value = [[rout_1, rout_1_v, rout_1_p], [rout_2, rout_2_v, rout_2_p], [rout_3, rout_3_v, rout_3_p], [rout_4, rout_4_v, rout_4_p]]
FourSquareSlide(prs, slide2_1_name, routs_value)
project1_name = 'Гарант в К0 (программа максимум)'
project1_value = ['Формирование команды проекта, проведены переговоры с:', 
                  'Руководством ВНИИПО, Жамойдиком С. проведены переговоры с начальником нормативно-технического отдела ДНПР МЧС России;', 
                  'Руководством УНК Пожарной безопасности в строительстве АГПС МЧС России; ', 
                  'Испытательной-пожарной лабораторией МЧС России по МО.', 
                  'Достигнуты предворительные договоренности о взаимодействии: в начале апреля ожидается получение проекта КП от АГПС.', 
                  'Совместно с Жамойдиком С. подготовлен проект плана действий, в апреле планируется представить план действий на обсуждение СБЕ.']
project1_limits = ['В ФГБУ ВНИИПО МЧС России происходит структурная перестройка, связанная со сменой руководства института, '
                   'это ведет к затягиванию решения вопросов организации работы, ',
                   'Прорабатывается вопрос определения АГПС МЧС России в качестве головной организации проекта. '
                   'Соответствующая позиция поддерживается руководством нормативно-технического отдела ДНПР МЧС России.']
part2_img1 = r'pics\2_1.jpg'

ProjectSlide(prs, project1_name, project1_value, project1_limits, img_link=part2_img1)

project2_name = 'Гарант в К0 (программа минимум)'
project2_value = ['Проведены испытания:',
                  'Решения покрытия системы ТН-Кровля Гарант, совместно с балками. Получен результат К1. Дальнейшие испытания в данном направлении остановлены ' \
                  'в связи с отказом Жамойдика С. признавать результаты таких испытаний в качестве "выполнения условий, определенных Войловым Е."', \
                  'Готовятся испытания:',
                  '2-х решений с введением систем огнезащиты в состав конструкции ТН-Кровля Гарант']
project2_limits = ['Возможность "публичного" прохождения испытаний на К0 ограничивается объективными факторами.']
part2_img2 = r'pics\2_2.jpg'

ProjectSlide(prs, project2_name, project2_value, project2_limits, img_link=part2_img2)

project3_name = 'Развитие пожарной комплектации'
project3_value = ['Проведены испытания 2-х решений огнестойких проходок водосточных воронок в систем ТН-Кровля Гарант. Подтверждено достижение огнестойкости не менее E15;',
                  'Информация о результатах поисковых испытаний передана отделу "Комплектации" для оценки коммерческой составляющей возможности введения решений огнезащиты воронок.']
project3_limits = ['Необходима оценка коммерческой составляющей введения огнестойких решений в комплектацию СБЕ.']
part2_img3 = r'pics\2_3.jpg'

ProjectSlide(prs, project3_name, project3_value, project3_limits, img_link=part2_img3)

project4_name = 'Развитие лаборатории пожарных испытаний'
project4_value = ['Запущен процесс проектирования ремонта помещений и вентиляции ЛПИ под размещение методов КП0 (ГОСТ Р 56026) и РП (ГОСТ Р 51032);',
                  'С 1 января 2025 г. запущена (в тестовом режиме) обновленная система приема и обработки заявок на проведения испытаний;', 
                  'Начато формирование общей базы данных испытаний проводимых в интересах СБЕ (в том числе проводимых во внешних лабораториях), а также архива документации формируемой в рамках работы лаборатории и направления ПБ;',
                  'Ведется разработка программных средств автоматизации процесса обработки результатов испытаний, формирования отчетной документации.']
project4_limits = ['Работа лаборатории харатеризуется низкой устойчивостью к "перегрузкам", в феврале, в связи с нахождением руководителя лаборатории на больничном, в '
                    'очередной раз отмечалось формирование задержек с выполнением заявок в интересах завода Logicroof',
                   'Необходимо введение в штате лаборатории одной дополнительной должности лаборантаи (испытателя), что, в том числе, позволит снизить сроки реализации '
                   'испытаний во внешних лабораториях, имеющих стратегическое значение для СБЕ. В настоящее время организация и участие в проведении испытаний'
                    'полностью находится в зоне компетенции Руководителя направления, привлечение Шои В. к осуществлению соответствующих работ ограничено, в связи '
                    'с необходимостью его нахождения в лаборатории.']
part2_img4 = r'pics\2_4.jpeg'

ProjectSlide(prs, project4_name, project4_value, project4_limits, img_link=part2_img4)

'''Завершаем 2 раздел'''


'''Начинаем 3 раздел: Лаборатория'''
PartTitleSlide(prs, 'Лаборатория пожарных испытаний', period_name, part3_title_pic_rout)

df_trecker = pd.read_excel(r'treck\treck.xlsx')
df_trecker = DateFormat(df_trecker, ['Дата начала', 'Дата завершения'])
df_trecker = SelectAndFilter(df_trecker, ['Ключ', 'Дата начала', 'Статус'], date_column='Дата начала', start_date=start_date, end_date=end_date)
df_trecker_activ = SelectAndFilter(df_trecker, ['Ключ', 'Дата начала', 'Статус'], column_name='Статус', filters=['В работе', 'Открыт', 'Требуется информация'])
report_column = ['Период', 'Количество заявок, всего', f'Количество не завершенных заявок на {end_date}', 'Количество испытаний в ЛПИ, всего','Количество испытаний во внешних лабораториях, всего', 'Количество испытаний в ЛПИ по заявкам завода ПВХ', 'Количество испытаний в ЛПИ по заявкам завода PIR', 
                 'Количество испытаний по заявкам Техслужбы', 'Количество испытаний "товарных" мембран: горючесть',
                 'Количество успешных испытаний "товарных" мембран: горючесть', 'Количество испытаний "товарных" мембран: воспламеняемость',
                 'Количество успешных испытаний "товарных" мембран: воспламеняемость', 'Количество испытаний "товарного" ПИР: горючесть', 'Количество успешных испытаний "товарного" ПИР: горючесть']
# print(df_trecker)

investigation_df = pd.read_excel(r'rep\report.xlsx')
invest_list = ['Предмет исследования', 'Номер заявки', 'Номер испытания', 'e-mail заказчика', 'ФИО заказчика', 'Наименование лаборатории', 'ЕКН', 'Название материала', 
               'Дата протокола', 'ФИО испытателя', 'Дата поступления образцов', 'Дата проведения испытания',  'Группа горючести по температуре ДГ', 'Соответствие', 
               'Группа горючести по повреждениям', 'Соответствие.1', 'Группа горючести по потере массы', 'Соответствие.2', 'Группа горючести по времени самостоятельного горения', 
               'Соответствие.3', 'Группа горючести по падению капель', 'Соответствие декларируемой группе горючести по падению капель', 'Итоговая оценка группы горючести', 
               'Общее соответствие', 'Дата регистрации результатов', 'Испытатель', 'Дата поступления образцов.1', 'Дата проведения испытания.1', 
                'Установленная группа воспламеняемости', 'Соответствие.4', 'Группа заказчиков']
investigation_df_temp = SelectAndFilter(investigation_df, invest_list, key_column='Наименование лаборатории', key_word='ЛПИ')
temp_df = investigation_df_temp.copy().reset_index(drop=True, level=0)
lpi_inc_set = set()
for i_x in range(0, len(investigation_df_temp)):        
    i_v = int(temp_df.at[i_x, 'Номер заявки'])
    lpi_inc_set.add(i_v)
# print(lpi_inc_set)

lpi_df = investigation_df[investigation_df['Номер заявки'].isin(lpi_inc_set)]
lpi_df = SelectAndFilter(lpi_df, invest_list)
lpi_comb_df = SelectAndFilter(lpi_df, invest_list, date_column='Дата протокола', start_date=start_date, end_date=end_date)
lpi_flam_df = SelectAndFilter(lpi_df, invest_list, date_column='Дата регистрации результатов', start_date=start_date, end_date=end_date)
ext_flam_df = SelectAndFilter(investigation_df, invest_list, date_column='Дата регистрации результатов', start_date=start_date, end_date=end_date)
lpi_comb_df_pvc = SelectAndFilter(lpi_comb_df, invest_list, key_column='Группа заказчиков', key_word='pvc')
lpi_flam_df_pvc = SelectAndFilter(lpi_flam_df, invest_list, key_column='Группа заказчиков', key_word='pvc')
lpi_comb_df_pir = SelectAndFilter(lpi_comb_df, invest_list, key_column='Группа заказчиков', key_word='pir')
lpi_flam_df_pir = SelectAndFilter(lpi_flam_df, invest_list, key_column='Группа заказчиков', key_word='pir')
lpi_comb_df_pvc_prod = SelectAndFilter(lpi_comb_df, invest_list, key_column='Название материала', key_word='мембрана')
lpi_flam_df_pvc_prod = SelectAndFilter(lpi_flam_df, invest_list, key_column='Название материала', key_word='мембрана')
lpi_comb_df_pir_prod = SelectAndFilter(lpi_comb_df, invest_list, key_column='Название материала', key_word='LOGICPIR')
lpi_comb_df_pvc_prod_match = SelectAndFilter(lpi_comb_df_pvc_prod, invest_list, column_name='Общее соответствие', filters=['Соответствует'])
lpi_flam_df_pvc_prod_match = SelectAndFilter(investigation_df, invest_list, column_name='Соответствие.4', filters=['Соответствует'])
lpi_comb_df_pir_prod_match = SelectAndFilter(lpi_comb_df_pir_prod, invest_list, column_name='Общее соответствие', filters=['Соответствует'])
# print(investigation_df)

# pvc_inv = SelectAndFilter(investigation_df, )

report_df = pd.DataFrame(columns=report_column)
report_df.at[0, "Период"] = period_name
report_df.at[0, "Количество заявок, всего"] = len(df_trecker)
report_df.at[0, f'Количество не завершенных заявок на {end_date}'] = len(df_trecker_activ)
report_df.at[0, 'Количество испытаний в ЛПИ, всего'] = len(lpi_comb_df) + len(lpi_flam_df)
report_df.at[0, 'Количество испытаний во внешних лабораториях, всего'] = len(ext_flam_df) - (len(lpi_comb_df) + len(lpi_flam_df))
report_df.at[0, 'Количество испытаний в ЛПИ по заявкам завода ПВХ'] = len(lpi_comb_df_pvc) + len(lpi_flam_df_pvc)
report_df.at[0, 'Количество испытаний в ЛПИ по заявкам завода PIR'] = len(lpi_comb_df_pir) + len(lpi_flam_df_pir)
report_df.at[0, 'Количество испытаний по заявкам Техслужбы'] = (report_df.at[0, 'Количество испытаний в ЛПИ, всего'] - 
                                                                report_df.at[0, 'Количество испытаний в ЛПИ по заявкам завода ПВХ'] - 
                                                                report_df.at[0, 'Количество испытаний в ЛПИ по заявкам завода PIR'])
report_df.at[0, 'Количество испытаний "товарных" мембран: горючесть'] = len(lpi_comb_df_pvc_prod)
report_df.at[0, 'Количество успешных испытаний "товарных" мембран: горючесть'] = len(lpi_comb_df_pvc_prod_match)
report_df.at[0, 'Количество испытаний "товарных" мембран: воспламеняемость'] = len(lpi_flam_df_pvc_prod)
report_df.at[0, 'Количество успешных испытаний "товарных" мембран: воспламеняемость'] = len(lpi_flam_df_pvc_prod_match)
report_df.at[0, 'Количество испытаний "товарного" ПИР: горючесть'] = len(lpi_comb_df_pir_prod)
report_df.at[0, 'Количество успешных испытаний "товарного" ПИР: горючесть'] = len(lpi_comb_df_pir_prod_match)
try:
    with pd.ExcelWriter(r'rep\report.xlsx', mode='a', engine='openpyxl') as writer:
        report_df.to_excel(writer, sheet_name='Производительность', index=False)
except:
    pass
report_to_plot = SelectAndFilter(report_df, report_column[1:])

## Создание графика
ax =report_to_plot.T.plot(  # Транспонируем DataFrame
 kind='bar',
 figsize=(11, 4.8),
 color='green',
 grid=True, 
 ylabel='Количество, ед.',
 legend=False # Отключаем автоматическую легенду
)

# Получаем названия категорий и высоты баров
categories = report_to_plot.columns
heights = report_to_plot.iloc[0].values

# Настраиваем шкалу Y
ax.set_ylim(0, max(heights) * 1.2) # Устанавливаем верхнюю границу с запасом
ax.yaxis.set_ticks(range(0, int(max(heights) * 1.2) + 1, 5)) # Устанавливаем шаг 5

# Размещаем названия категорий на графике
for i, (height, cat) in enumerate(zip(heights, categories)):
    ax.text(
        i - 0.35,                    # позиция по x
        1,         # позиция по y (немного выше бара)
        cat,                 # текст категории
        ha='center',         # выравнивание по горизонтали
        va='bottom',         # выравнивание по вертикали
        rotation=90,
        fontsize=5          # размер шрифта
    )

# Настраиваем остальные параметры графика
ax.set_xticks([])           # убираем метки по оси X
ax.set_xlabel('')           # убираем название оси X
plt.ylim(top=max(heights) * 1.2) # увеличиваем верхнее ограничение по Y


plt.savefig('chart.png', bbox_inches='tight', dpi=300)

slide_layout_3_1 = prs.slide_layouts[10]
slide_3_1 = prs.slides.add_slide(slide_layout_3_1)
CheckPlaceholders(slide_3_1)
image_placeholder_3_1 = slide_3_1.placeholders[13]
image_placeholder_3_1.insert_picture('chart.png')
title_3_1 = slide_3_1.shapes.title
title_3_1.text = 'ЛПИ: Сводные данные за I квартал 2025 г.'
title_3_1.text_frame.paragraphs[0].runs[0].font.size = Pt(28)


column_to_pic = ['Соответствие', 'Соответствие.1', 'Соответствие.2', 'Соответствие.3',
                 'Соответствие декларируемой группе горючести по падению капель', 'Общее соответствие']
column_to_pic_new = ['Температура дымовых газов', 'Длина повреждений',  'Потеря массы',
                         'Время самостоятельного горения', 'Падение горящих капель', 'Общая оценка']

column_to_pic_rename = dict(zip(column_to_pic, column_to_pic_new))
# print(column_to_pic_rename)
df_to_pic_pvc = lpi_comb_df_pvc_prod[column_to_pic]
df_to_pic_pvc = df_to_pic_pvc.rename(columns=column_to_pic_rename)
# print(df_to_pic)
PicatorVtoroy(df_to_pic_pvc, column_to_pic_new,  'LOGIROOF\n параметры соответствия', 'logicroof', 'Соответствует', 'Не соответствует', 'н/у', width=11, hight=4)


slide_3_2 = prs.slides.add_slide(slide_layout_3_1)
# CheckPlaceholders(slide_3_1)
image_placeholder_3_2 = slide_3_2.placeholders[13]
image_placeholder_3_2.insert_picture('logicroof.jpg')
title_3_2 = slide_3_2.shapes.title
title_3_2.text = 'LOGICROOF параметры (не) соответствия: горючесть.'
title_3_2.text_frame.paragraphs[0].runs[0].font.size = Pt(28)



df_to_pic = lpi_comb_df_pir[column_to_pic]
df_to_pic = df_to_pic.rename(columns=column_to_pic_rename)
# print(df_to_pic)
PicatorVtoroy(df_to_pic, column_to_pic_new,  'LOGICPIR\n параметры соответствия', 'logicpir', 'Соответствует', 'Не соответствует', width=11, hight=4)


slide_3_3 = prs.slides.add_slide(slide_layout_3_1)
# CheckPlaceholders(slide_3_1)
image_placeholder_3_3 = slide_3_3.placeholders[13]
image_placeholder_3_3.insert_picture('logicpir.jpg')
title_3_3 = slide_3_3.shapes.title
title_3_3.text = 'LOGICPIR параметры (не) соответствия: горючесть.'
title_3_3.text_frame.paragraphs[0].runs[0].font.size = Pt(28)

# # print(report_df)
# report_df.to_excel('temp.xlsx')
# os.startfile('temp.xlsx')

'''Заканчиваем 3 раздел'''


prs.save(report_file)

os.startfile(report_file)