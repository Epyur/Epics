import os
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from func import *


def CheckPlaceholders(slide_name):
    for placeholder in slide_name.placeholders:
        print(placeholder.placeholder_format.idx, placeholder.placeholder_format.type)

def FormatText(place_holder, size=None, color=None, bold=None, italic=None, align=None):
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
    # Форматирование текста
    text_frame = title_first_part.text_frame
    paragraph = text_frame.paragraphs[0]
    run = paragraph.runs[0]

    if size is not None:
        # Размер шрифта
        run.font.size = Pt(size)

    if color is not None:
        # Цвет шрифта
        run.font.color.rgb = RGBColor(0, 0, 255)  # синий цвет

    # Начертание
    if bold is not None:
        run.font.bold = True
    if italic is not None:    
        run.font.italic = True

    if align == 'L':
        # Выравнивание
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    if align == 'R':
        # Выравнивание
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
    if align == 'C':
        # Выравнивание
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    if align == 'J':
        # Выравнивание
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.JUSTIFY

# def TransparencyImg(shape, alpha):
    
#     ts = shape.fill._xPr.solidFill
#     sF = ts.get
#     sE = SubElement(sF, 'a:alpha', val=str(alpha))


start_pres = os.path.abspath(r'C:\Users\polishchuk\Documents\temple\temple.pptx')

prs = Presentation(start_pres)
# prs = prs.slide_layouts[19]

'''Формируем титульный слайд'''

slide_layout = prs.slide_layouts[2]  # Заголовок и контент
slide_title = prs.slides.add_slide(slide_layout)

# CheckPlaceholders(slide_title)

title_first = slide_title.shapes.title
title_first.text = 'Направление пожарная безопасность'
title_first.text_frame.paragraphs[0].runs[0].font.size = Pt(28)

subtitle_1 = slide_title.placeholders[13]
subtitle_1.text = 'I квартал'

subtitle_0 = slide_title.placeholders[12]
subtitle_0.text = '25'

image_placeholder = slide_title.placeholders[11]
image_placeholder.insert_picture(r'Gen_4rep\pics\title.png')

'''Закончили формировать титульный слайд'''

'''Начинаем 1 раздел: мероприятия'''
slide_layout_part = prs.slide_layouts[4]
slide_part1_title = prs.slides.add_slide(slide_layout_part)
# CheckPlaceholders(slide_part1_title)

title_first_part = slide_part1_title.placeholders[10]
title_first_part.text = 'Техническая поддержка и мероприятия'



subtitle_2 = slide_part1_title.placeholders[11]
subtitle_2.text = 'I квартал'

FormatText(title_first_part, align='C')
image_placeholder_2 = slide_part1_title.placeholders[12]
image_placeholder_2.insert_picture(r'Gen_4rep\pics\part1.png')


text_1 = ['1. Собрание СБЕ']
slide_1_1_layouts = prs.slide_layouts[11]
slide_1_1 = prs.slides.add_slide(slide_1_1_layouts)
# CheckPlaceholders(slide_1_1)
title_1_1 = slide_1_1.shapes.title
title_1_1.text = 'Текущая работа в I квартале 2025 года'
title_1_1.text_frame.paragraphs[0].runs[0].font.size = Pt(28)

body_1_1 = slide_1_1.placeholders[13]
body_1_1_tf = body_1_1.text_frame
body_1_1_tf.clear()
body_1_1_pg = body_1_1_tf.add_paragraph()
body_1_1_pg.add_run().text = '18-20 февраля: Участие в собрании СБЕ; \n ' \
'январь - февраль: Участие в совещаниях по Саларьево (МИП); \n ' \
'январь - февраль: Участие в проведении испытаний мембран в ЦЭИИС, в рамках работы по объекту "Саларьево"; \n ' \
'Общая техническая поддержка (подготовка писем, проведение консультаций);\n' \
'Получен сертификат на систему ТН-Кровля Гарант в системе НСОПБ с размещением информации о сертификате в открытых реестрах системы сертификации; \n' \
'Получен протокол Г1 на PIRCRYO, по результатам испытаний во ВНИИПО; \n ' \
'30-31 января: Участие в конференции "Безопасная Арктика" в г. Оренбург;\n' \
'20 марта: Участие в конференции "Ройтмановские чтения" в АГПС (г. Москва) совместно с Жамойдиком С.' \
'\n Пожар в клину: выезд на объект (снятие сомнений по вопросу применения системы ТН-Кровля Гарант)'

'''Завершаем 1 раздел'''

'''Начинаем 2 раздел: НИОКР и проектная деятельность'''
slide_layout_2 = prs.slide_layouts[4]
slide_part2_title = prs.slides.add_slide(slide_layout_2)
# CheckPlaceholders(slide_part2_title)

title_first_2_0 = slide_part2_title.placeholders[10]
title_first_2_0.text = 'НИОКРы и проектная деятельность'
FormatText(title_first_2_0, align='C')

subtitle_2_0 = slide_part2_title.placeholders[11]
subtitle_2_0.text = 'I квартал'

image_placeholder_2_0 = slide_part2_title.placeholders[12]
image_placeholder_2_0.insert_picture(r'Gen_4rep\pics\part2.png')

# Part 2_1
slide_layout_2_1 = prs.slide_layouts[30]
slide_2_1 = prs.slides.add_slide(slide_layout_2_1)
# CheckPlaceholders(slide_2_1)

title_2_1 = slide_2_1.shapes.title
title_2_1.text = 'Основные направления движения'
title_2_1.text_frame.paragraphs[0].runs[0].font.size = Pt(28)

item_2_1_1 = slide_2_1.placeholders[13]
item_2_1_1.text = 'Гарант в К0 (программа максимум)'
item_2_1_1.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
image_placeholder_2_1_1 = slide_2_1.placeholders[28]
image_placeholder_2_1_1.insert_picture(r'Gen_4rep\pics\10.png')

item_2_1_1_aim = slide_2_1.placeholders[26]
item_2_1_1_aim.text = 'Задача: снять внешние риски и ограничения для продвижения системы ТН-Кровля Гарант'
item_2_1_1_aim.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
item_2_1_1_aim.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 0, 0)


item_2_1_2 = slide_2_1.placeholders[14]
item_2_1_2.text = 'Гарант в К0 (программа минимум)'
item_2_1_2.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
image_placeholder_2_1_2 = slide_2_1.placeholders[22]
image_placeholder_2_1_2.insert_picture(r'Gen_4rep\pics\50.png')
item_2_1_2_aim = slide_2_1.placeholders[25]
item_2_1_2_aim.text = 'Задача: снять внутрикорпоративные ограничения на продвижение системы ТН-Кровля Гарант'
item_2_1_2_aim.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
item_2_1_2_aim.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 0, 0)

item_2_1_3 = slide_2_1.placeholders[16]
item_2_1_3.text = 'Развитие пожарной комплектации'
item_2_1_3.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
image_placeholder_2_1_3 = slide_2_1.placeholders[29]
image_placeholder_2_1_3.insert_picture(r'Gen_4rep\pics\75.png')
item_2_1_3_aim = slide_2_1.placeholders[27]
item_2_1_3_aim.text = 'Задачи: \n Развитие комплектности предложений СБЕ;' \
' \n Формирование конкурентных преимуществ для системных решений СБЕ;'
item_2_1_3_aim.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
item_2_1_3_aim.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 0, 0)
item_2_1_3_aim.text_frame.paragraphs[1].runs[0].font.size = Pt(14)
item_2_1_3_aim.text_frame.paragraphs[1].runs[0].font.color.rgb = RGBColor(255, 0, 50)
item_2_1_3_aim.text_frame.paragraphs[2].runs[0].font.size = Pt(14)
item_2_1_3_aim.text_frame.paragraphs[2].runs[0].font.color.rgb = RGBColor(255, 0, 100)

item_2_1_4 = slide_2_1.placeholders[15]
item_2_1_4.text = 'Развитие ЛПИ и исследований материалов СБЕ'
item_2_1_4.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
image_placeholder_2_1_4 = slide_2_1.placeholders[30]
image_placeholder_2_1_4.insert_picture(r'Gen_4rep\pics\50.png')
item_2_1_4_aim = slide_2_1.placeholders[21]
item_2_1_4_aim.text = 'Задачи: \n Проведение модернизации исследовательской базы лаборатории;' \
' \n Выстраивание процессов контроля работы лаборатории, включая формирование базы результатов проводимых исследований;' \
'\n Проведение исследований показателей воспламеняемости и распространения пламени по поверхности ПВХ мембран.'
item_2_1_4_aim.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
item_2_1_4_aim.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 0, 0)
item_2_1_4_aim.text_frame.paragraphs[1].runs[0].font.size = Pt(13)
item_2_1_4_aim.text_frame.paragraphs[1].runs[0].font.color.rgb = RGBColor(255, 0, 50)
item_2_1_4_aim.text_frame.paragraphs[2].runs[0].font.size = Pt(13)
item_2_1_4_aim.text_frame.paragraphs[2].runs[0].font.color.rgb = RGBColor(255, 0, 100)
item_2_1_4_aim.text_frame.paragraphs[3].runs[0].font.size = Pt(13)
item_2_1_4_aim.text_frame.paragraphs[3].runs[0].font.color.rgb = RGBColor(255, 0, 150)

slide_layout_2_2 = prs.slide_layouts[32]
slide_2_2 = prs.slides.add_slide(slide_layout_2_2)
# CheckPlaceholders(slide_2_2)
image_placeholder_2_2_1 = slide_2_2.placeholders[13]
image_placeholder_2_2_1.insert_picture(r'Gen_4rep\pics\2_2.jpg')
title_2_2 = slide_2_2.shapes.title
title_2_2.text = 'Гарант в К0 (программа максимум)'
title_2_2.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
item_2_2_aim = slide_2_2.placeholders[14]
item_2_2_aim.text = 'Действия и намерения:' \
' \n Формирование команды проекта, проведены переговоры с:' \
'\n         Руководством ВНИИПО, Жамойдиком С. проведены переговоры с начальником нормативно-технического отдела ДНПР МЧС России;' \
'\n         Руководством УНК Пожарной безопасности в строительстве; ' \
'\n         Санкт-Петербургским университетом ГПС МЧС России;' \
'\n         Испытательной-пожарной лабораторией МЧС России по МО.' \
'\n Достигнуты предворительные договоренности о взаимодействии.' \
'\n Совместно с Жамойдиком С. подготовлен проект плана действий;' \
'\n В апреле планируется представить план действий на обсуждение СБЕ.'
item_2_2_aim.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
item_2_2_aim.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 255)
for c_l in range(1, 9):
    item_2_2_aim.text_frame.paragraphs[c_l].runs[0].font.size = Pt(10)
    item_2_2_aim.text_frame.paragraphs[c_l].runs[0].font.color.rgb = RGBColor(47, 79, 79)

item_2_2_aim = slide_2_2.placeholders[15]
item_2_2_aim.text = 'Ограничения:' \
' \n В ФГБУ ВНИИПО МЧС России происходит структурная перестройка, связанная со сменой руководства института, это ведет к затягиванию решения вопросов организации работы, ' \
'прорабатывается вопрос определения АГПС МЧС России в качестве головной организации проекта. Соответствующая позиция поддерживается руководством нормативно-технического отдела ДНПР' \
'МЧС России.'
item_2_2_aim.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
item_2_2_aim.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 255)
item_2_2_aim.text_frame.paragraphs[1].runs[0].font.size = Pt(12)
item_2_2_aim.text_frame.paragraphs[1].runs[0].font.color.rgb = RGBColor(47, 79, 79)


slide_layout_2_3 = prs.slide_layouts[32]
slide_2_3 = prs.slides.add_slide(slide_layout_2_3)
title_2_3 = slide_2_3.shapes.title
title_2_3.text = 'Гарант в К0 (программа минимум)'
title_2_3.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
item_2_3_aim = slide_2_3.placeholders[14]
item_2_3_aim.text = 'Действия и намерения:' \
' \n Проведены испытания:' \
'\n         Решения покрытия системы ТН-Кровля Гарант, совместно с балками. Получен результат К1. Дальнейшие испытания в данном направлении остановлены ' \
'в связи с отказом Жамойдика С. признавать результаты таких испытаний в качестве "выполнения условий, определенных Войловым Е."' \
' \n Готовятся испытания:' \
'\n     2-х решений с введением систем огнезащиты в состав конструкции ТН-Кровля Гарант'
item_2_3_aim.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
item_2_3_aim.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 255)
for c_l in range(1, 5):
    item_2_3_aim.text_frame.paragraphs[c_l].runs[0].font.size = Pt(10)
    item_2_3_aim.text_frame.paragraphs[c_l].runs[0].font.color.rgb = RGBColor(47, 79, 79)

item_2_3_aim = slide_2_3.placeholders[15]
item_2_3_aim.text = 'Ограничения:' \
' \n Возможность "публичного" прохождения испытаний на К0 ограничивается объективными факторами.'
item_2_3_aim.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
item_2_3_aim.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 255)
item_2_3_aim.text_frame.paragraphs[1].runs[0].font.size = Pt(12)
item_2_3_aim.text_frame.paragraphs[1].runs[0].font.color.rgb = RGBColor(47, 79, 79)
image_placeholder_2_3 = slide_2_3.placeholders[13]
image_placeholder_2_3.insert_picture(r'Gen_4rep\pics\2_3.jpg')


slide_layout_2_4 = prs.slide_layouts[32]
slide_2_4 = prs.slides.add_slide(slide_layout_2_4)
title_2_4 = slide_2_4.shapes.title
title_2_4.text = 'Развитие пожарной комплектации'
title_2_4.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
item_2_4_aim = slide_2_4.placeholders[14]
item_2_4_aim.text = 'Действия и намерения:' \
' \n Проведены испытания 2-х решений огнестойких проходок водосточных воронок в систем ТН-Кровля Гарант. Подтверждено достижение огнестойкости не менее E15;' \
'\n Информация о результатах поисковых испытаний передана отделу "Комплектации" для оценки коммерческой составляющей возможности введения решений огнезащиты воронок.'
item_2_4_aim.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
item_2_4_aim.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 255)
for c_l in range(1, 3):
    item_2_4_aim.text_frame.paragraphs[c_l].runs[0].font.size = Pt(11)
    item_2_4_aim.text_frame.paragraphs[c_l].runs[0].font.color.rgb = RGBColor(47, 79, 79)

item_2_4_aim = slide_2_4.placeholders[15]
item_2_4_aim.text = 'Ограничения:' \
' \n Необходима оценка коммерческой составляющей введения огнестойких решений в комплектацию СБЕ.'
item_2_4_aim.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
item_2_4_aim.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 255)
item_2_4_aim.text_frame.paragraphs[1].runs[0].font.size = Pt(12)
item_2_4_aim.text_frame.paragraphs[1].runs[0].font.color.rgb = RGBColor(47, 79, 79)
image_placeholder_2_4 = slide_2_4.placeholders[13]
image_placeholder_2_4.insert_picture(r'Gen_4rep\pics\2_4.jpg')

slide_layout_2_5 = prs.slide_layouts[32]
slide_2_5 = prs.slides.add_slide(slide_layout_2_4)
title_2_5 = slide_2_5.shapes.title
title_2_5.text = 'Развитие лаборатории пожарных испытаний'
title_2_5.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
item_2_5_aim = slide_2_5.placeholders[14]
item_2_5_aim.text = 'Действия и намерения:' \
' \n Запущен процесс проектирования ремонта помещений и вентиляции ЛПИ под размещение методов КП0 (ГОСТ Р 56026) и РП (ГОСТ Р 51032);' \
'\n С 1 января 2025 г. запущена (в тестовом режиме) обновленная система приема и обработки заявок на проведения испытаний;' \
'\n Начато формирование общей базы данных испытаний проводимых в интересах СБЕ (в том числе проводимых во внешних лабораториях), а также архива документации формируемой в рамках работы лаборатории и направления ПБ;' \
'\n Ведется разработка программных средств автоматизации процесса обработки результатов испытаний, формирования отчетной документации.'

item_2_5_aim.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
item_2_5_aim.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 255)
for c_l in range(1, 5):
    item_2_5_aim.text_frame.paragraphs[c_l].runs[0].font.size = Pt(11)
    item_2_5_aim.text_frame.paragraphs[c_l].runs[0].font.color.rgb = RGBColor(47, 79, 79)

item_2_5_aim = slide_2_5.placeholders[15]
item_2_5_aim.text = 'Ограничения:' \
' \n Работа лаборатории харатеризуется низкой устойчивостью к "перегрузкам", в феврале, в связи с нахождением руководителя лаборатории на больничном, в очередной раз отмечалось формирование задержек с выполнением заявок в интересах завода Logicroof' \
'\n Необходимо введение в штате лаборатории одной дополнительной должности лаборантаи (испытателя), что, в том числе, позволит снизить сроки реализации испытаний во внешних лабораториях, имеющих стратегическое значение для СБЕ. В настоящее время организация и участие в проведении испытаний' \
'полностью находится в зоне компетенции Руководителя направления, привлечение Шои В. к осуществлению соответствующих работ ограничено, в связи с необходимостью его нахождения в лаборатории.'
item_2_5_aim.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
item_2_5_aim.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 255)
item_2_5_aim.text_frame.paragraphs[1].runs[0].font.size = Pt(12)
item_2_5_aim.text_frame.paragraphs[1].runs[0].font.color.rgb = RGBColor(47, 79, 79)
item_2_5_aim.text_frame.paragraphs[2].runs[0].font.size = Pt(12)
item_2_5_aim.text_frame.paragraphs[2].runs[0].font.color.rgb = RGBColor(47, 79, 79)
image_placeholder_2_5 = slide_2_5.placeholders[13]
image_placeholder_2_5.insert_picture(r'Gen_4rep\pics\2_5.jpeg')

'''Завершаем 2 раздел'''


'''Начинаем 3 раздел: Лаборатория'''

slide_part3_title = prs.slides.add_slide(slide_layout_2)
# CheckPlaceholders(slide_part2_title)

title_first_3_0 = slide_part3_title.placeholders[10]
title_first_3_0.text = 'Лаборатория пожарных испытаний'
FormatText(title_first_3_0, align='C')

subtitle_3_0 = slide_part3_title.placeholders[11]
subtitle_3_0.text = 'I квартал'

image_placeholder_3_0 = slide_part3_title.placeholders[12]
image_placeholder_3_0.insert_picture(r'Gen_4rep\pics\part3.jpeg')


df_trecker = pd.read_excel(r'Gen_4rep\treck\treck.xlsx')
start_date = '01.01.2025'
end_date = '31.03.2025'
df_trecker = DateFormat(df_trecker, ['Обновлено', 'Дата начала', 'Дата завершения'])
df_trecker = SelectAndFilter(df_trecker, ['Ключ', 'Дата начала'], date_column='Дата начала', start_date=start_date, end_date=end_date)
report_column = ['Количество заявок, всего', 'Количество испытаний в ЛПИ, всего', 'Количество испытаний в ЛПИ по заявкам завода ПВХ', 'Количество испытаний в ЛПИ по заявкам завода PIR', 
                 'Количество испытаний по заявкам Техслужбы', 'Количество испытаний "товарных" мембран: горючесть',
                 'Количество успешных испытаний "товарных" мембран: горючесть', 'Количество испытаний "товарных" мембран: воспламеняемость',
                 'Количество успешных испытаний "товарных" мембран: воспламеняемость', 'Количество испытаний "товарного" ПИР: горючесть', 'Количество успешных испытаний "товарного" ПИР: горючесть']
# print(df_trecker)

investigation_df = pd.read_excel(r'Gen_4rep\rep\report.xlsx')
invest_list = ['Предмет исследования', 'Номер заявки', 'Номер испытания', 'e-mail заказчика', 'ФИО заказчика', 'Наименование лаборатории', 'ЕКН', 'Название материала', 
               'Дата протокола', 'ФИО испытателя', 'Дата поступления образцов', 'Дата проведения испытания',  'Группа горючести по температуре ДГ', 'Соответствие', 
               'Группа горючести по повреждениям', 'Соответствие.1', 'Группа горючести по потере массы', 'Соответствие.2', 'Группа горючести по времени самостоятельного горения', 
               'Соответствие.3', 'Группа горючести по падению капель', 'Соответствие декларируемой группе горючести по падению капель', 'Итоговая оценка группы горючести', 
               'Общее соответствие', 'Дата регистрации результатов', 'Испытатель', 'Дата поступления образцов.1', 'Дата проведения испытания.1', 
                'Установленная группа воспламеняемости', 'Соответствие.4', 'Группа заказчиков']
investigation_df_temp = SelectAndFilter(investigation_df, invest_list, key_column='Наименование лаборатории', key_word='ЛПИ')
temp_df = investigation_df_temp.copy().reset_index(drop=True, level=0)
lpi_inc_set = set()
for i_x in range(0, len(investigation_df_temp)):        
    i_v = int(temp_df.at[i_x, 'Номер заявки'])
    lpi_inc_set.add(i_v)
# print(lpi_inc_set)

lpi_df = investigation_df[investigation_df['Номер заявки'].isin(lpi_inc_set)]
lpi_df = SelectAndFilter(lpi_df, invest_list)
lpi_comb_df = SelectAndFilter(lpi_df, invest_list, date_column='Дата протокола', start_date=start_date, end_date=end_date)
lpi_flam_df = SelectAndFilter(lpi_df, invest_list, date_column='Дата регистрации результатов', start_date=start_date, end_date=end_date)
lpi_comb_df_pvc = SelectAndFilter(lpi_comb_df, invest_list, key_column='Группа заказчиков', key_word='pvc')
lpi_flam_df_pvc = SelectAndFilter(lpi_flam_df, invest_list, key_column='Группа заказчиков', key_word='pvc')
lpi_comb_df_pir = SelectAndFilter(lpi_comb_df, invest_list, key_column='Группа заказчиков', key_word='pir')
lpi_flam_df_pir = SelectAndFilter(lpi_flam_df, invest_list, key_column='Группа заказчиков', key_word='pir')
lpi_comb_df_pvc_prod = SelectAndFilter(lpi_comb_df, invest_list, key_column='Название материала', key_word='мембрана')
lpi_flam_df_pvc_prod = SelectAndFilter(lpi_flam_df, invest_list, key_column='Название материала', key_word='мембрана')
lpi_comb_df_pir_prod = SelectAndFilter(lpi_comb_df, invest_list, key_column='Название материала', key_word='LOGICPIR')
lpi_comb_df_pvc_prod_match = SelectAndFilter(lpi_comb_df_pvc_prod, invest_list, column_name='Общее соответствие', filters=['Соответствует'])
lpi_flam_df_pvc_prod_match = SelectAndFilter(investigation_df, invest_list, column_name='Соответствие.4', filters=['Соответствует'])
lpi_comb_df_pir_prod_match = SelectAndFilter(lpi_comb_df_pir_prod, invest_list, column_name='Общее соответствие', filters=['Соответствует'])
# print(investigation_df)

# pvc_inv = SelectAndFilter(investigation_df, )

report_df = pd.DataFrame(columns=report_column)
report_df.at[0, "Количество заявок, всего"] = len(df_trecker)
report_df.at[0, 'Количество испытаний в ЛПИ, всего'] = len(lpi_comb_df) + len(lpi_flam_df)
report_df.at[0, 'Количество испытаний в ЛПИ по заявкам завода ПВХ'] = len(lpi_comb_df_pvc) + len(lpi_flam_df_pvc)
report_df.at[0, 'Количество испытаний в ЛПИ по заявкам завода PIR'] = len(lpi_comb_df_pir) + len(lpi_flam_df_pir)
report_df.at[0, 'Количество испытаний по заявкам Техслужбы'] = (report_df.at[0, 'Количество испытаний в ЛПИ, всего'] - 
                                                                report_df.at[0, 'Количество испытаний в ЛПИ по заявкам завода ПВХ'] - 
                                                                report_df.at[0, 'Количество испытаний в ЛПИ по заявкам завода PIR'])
report_df.at[0, 'Количество испытаний "товарных" мембран: горючесть'] = len(lpi_comb_df_pvc_prod)
report_df.at[0, 'Количество успешных испытаний "товарных" мембран: горючесть'] = len(lpi_comb_df_pvc_prod_match)
report_df.at[0, 'Количество испытаний "товарных" мембран: воспламеняемость'] = len(lpi_flam_df_pvc_prod)
report_df.at[0, 'Количество успешных испытаний "товарных" мембран: воспламеняемость'] = len(lpi_flam_df_pvc_prod_match)
report_df.at[0, 'Количество испытаний "товарного" ПИР: горючесть'] = len(lpi_comb_df_pir_prod)
report_df.at[0, 'Количество успешных испытаний "товарного" ПИР: горючесть'] = len(lpi_comb_df_pir_prod_match)

## Создание графика
ax =report_df.T.plot(  # Транспонируем DataFrame
 kind='bar',
 figsize=(11, 4.8),
 color='green',
 grid=True, 
 ylabel='Значения',
 legend=False # Отключаем автоматическую легенду
)

# Получаем названия категорий и высоты баров
categories = report_df.columns
heights = report_df.iloc[0].values

# Настраиваем шкалу Y
ax.set_ylim(0, max(heights) * 1.2) # Устанавливаем верхнюю границу с запасом
ax.yaxis.set_ticks(range(0, int(max(heights) * 1.2) + 1, 5)) # Устанавливаем шаг 5

# Размещаем названия категорий на графике
for i, (height, cat) in enumerate(zip(heights, categories)):
    ax.text(
        i - 0.35,                    # позиция по x
        1,         # позиция по y (немного выше бара)
        cat,                 # текст категории
        ha='center',         # выравнивание по горизонтали
        va='bottom',         # выравнивание по вертикали
        rotation=90,
        fontsize=5          # размер шрифта
    )

# Настраиваем остальные параметры графика
ax.set_xticks([])           # убираем метки по оси X
ax.set_xlabel('')           # убираем название оси X
plt.ylim(top=max(heights) * 1.2) # увеличиваем верхнее ограничение по Y


plt.savefig('chart.png', bbox_inches='tight', dpi=300)

slide_layout_3_1 = prs.slide_layouts[10]
slide_3_1 = prs.slides.add_slide(slide_layout_3_1)
CheckPlaceholders(slide_3_1)
image_placeholder_3_1 = slide_3_1.placeholders[13]
image_placeholder_3_1.insert_picture('chart.png')
title_3_1 = slide_3_1.shapes.title
title_3_1.text = 'ЛПИ: Сводные данные за I квартал 2025 г.'
title_3_1.text_frame.paragraphs[0].runs[0].font.size = Pt(28)


column_to_pic = ['Соответствие', 'Соответствие.1', 'Соответствие.2', 'Соответствие.3',
                 'Соответствие декларируемой группе горючести по падению капель', 'Общее соответствие']
column_to_pic_new = ['Температура дымовых газов', 'Длина повреждений',  'Потеря массы',
                         'Время самостоятельного горения', 'Падение горящих капель', 'Общая оценка']

column_to_pic_rename = dict(zip(column_to_pic, column_to_pic_new))
# print(column_to_pic_rename)
df_to_pic_pvc = lpi_comb_df_pvc_prod[column_to_pic]
df_to_pic_pvc = df_to_pic_pvc.rename(columns=column_to_pic_rename)
# print(df_to_pic)
PicatorVtoroy(df_to_pic_pvc, column_to_pic_new,  'LOGIROOF\n параметры соответствия', 'logicroof', 'Соответствует', 'Не соответствует', 'н/у', width=11, hight=4)


slide_3_2 = prs.slides.add_slide(slide_layout_3_1)
# CheckPlaceholders(slide_3_1)
image_placeholder_3_2 = slide_3_2.placeholders[13]
image_placeholder_3_2.insert_picture('logicroof.jpg')
title_3_2 = slide_3_2.shapes.title
title_3_2.text = 'LOGICROOF параметры (не) соответствия: горючесть.'
title_3_2.text_frame.paragraphs[0].runs[0].font.size = Pt(28)



df_to_pic = lpi_comb_df_pir[column_to_pic]
df_to_pic = df_to_pic.rename(columns=column_to_pic_rename)
# print(df_to_pic)
PicatorVtoroy(df_to_pic, column_to_pic_new,  'LOGICPIR\n параметры соответствия', 'logicpir', 'Соответствует', 'Не соответствует', width=11, hight=4)


slide_3_3 = prs.slides.add_slide(slide_layout_3_1)
# CheckPlaceholders(slide_3_1)
image_placeholder_3_3 = slide_3_3.placeholders[13]
image_placeholder_3_3.insert_picture('logicpir.jpg')
title_3_3 = slide_3_3.shapes.title
title_3_3.text = 'LOGICPIR параметры (не) соответствия: горючесть.'
title_3_3.text_frame.paragraphs[0].runs[0].font.size = Pt(28)

# # print(report_df)
# report_df.to_excel('temp.xlsx')
# os.startfile('temp.xlsx')

'''Заканчиваем 3 раздел'''


prs.save('just_to_see.pptx')

os.startfile('just_to_see.pptx')