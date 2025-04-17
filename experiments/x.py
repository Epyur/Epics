import aspose.slides as slides
import os
from pptx import Presentation

trace_f = os.path.abspath(r'C:\Users\epyur\OneDrive\Документы\temple\temple.pptx')

with slides.Presentation(trace_f) as prc:
    print(list(prc.slides))
    for i in [40, 41, 42, 43, 44, 45]:
        slide = prc.slides[i]
        prc.slides.remove(slide)

    prc.save('new_one.pptx', slides.export.SaveFormat.PPTX)

















def remove_text(presentation_path, search_text):
    # Открываем презентацию
    prs = Presentation(presentation_path)

    # Проходим по всем слайдам
    for slide in prs.slides:
        # Проходим по всем текстовым фреймам на слайде
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame

                # Проходим по всем абзацам в текстовом фрейме
                for paragraph in text_frame.paragraphs:
                    # Проходим по всем runs (частям текста) в абзаце
                    for run in paragraph.runs:
                        # Если текст содержит искомое слово
                        if search_text in run.text:
                            # Заменяем текст на пустой
                            run.text = run.text.replace(search_text, "")

    # Сохраняем измененную презентацию
    prs.save(presentation_path)

remove_text('new_one.pptx', 'Evaluation only.')
remove_text('new_one.pptx','Created with Aspose.Slides for Python via .NET 25.2.')
remove_text('new_one.pptx','Copyright 2004-2025Aspose Pty Ltd.')

os.startfile('new_one.pptx')